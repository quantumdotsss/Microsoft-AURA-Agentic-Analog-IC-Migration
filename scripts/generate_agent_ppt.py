#!/usr/bin/env python3
from __future__ import annotations

import subprocess
import textwrap
from datetime import datetime
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parents[1]
OUT = BASE / "AURA_Agentic_Framework_Yunbo.pptx"
ASSETS = BASE / "presentation_assets"
ASSETS.mkdir(parents=True, exist_ok=True)

SLIDE_W = 13.333333
SLIDE_H = 7.5
FOOTER_H = 0.72
BLUE = RGBColor(0, 102, 161)
UCI_BLUE = RGBColor(0, 93, 148)
GOLD = RGBColor(255, 210, 0)
DARK = RGBColor(25, 25, 25)
MS_RED = RGBColor(242, 80, 34)
MS_GREEN = RGBColor(127, 186, 0)
MS_BLUE = RGBColor(0, 164, 239)
MS_YELLOW = RGBColor(255, 185, 0)
GREEN = RGBColor(46, 125, 50)
ORANGE = RGBColor(230, 81, 0)
MAGENTA = RGBColor(194, 24, 91)
LIGHT_BLUE = RGBColor(225, 245, 254)
LIGHT_GREEN = RGBColor(232, 245, 233)
LIGHT_ORANGE = RGBColor(255, 243, 224)
LIGHT_MAGENTA = RGBColor(252, 228, 236)
GRAY = RGBColor(95, 95, 95)


def font(name: str = "DejaVuSans.ttf", size: int = 24) -> ImageFont.FreeTypeFont:
    candidates = [
        f"/usr/share/fonts/truetype/dejavu/{name}",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        p = Path(candidate)
        if p.exists():
            return ImageFont.truetype(str(p), size=size)
    return ImageFont.load_default()


def code_image(text: str, out: Path, width: int = 1100, font_size: int = 22) -> Path:
    mono = font("DejaVuSansMono.ttf", font_size)
    line_h = int(font_size * 1.45)
    lines = text.splitlines()
    height = max(160, 34 + line_h * len(lines))
    img = Image.new("RGB", (width, height), (43, 7, 33))
    draw = ImageDraw.Draw(img)
    y = 18
    for line in lines:
        draw.text((22, y), line[:120], font=mono, fill=(235, 235, 235))
        y += line_h
    img.save(out)
    return out


def code_panel(text: str, out: Path, width: int = 1180, height: int = 520, font_size: int = 17) -> Path:
    mono = font("DejaVuSansMono.ttf", font_size)
    line_h = int(font_size * 1.42)
    img = Image.new("RGB", (width, height), (43, 7, 33))
    draw = ImageDraw.Draw(img)
    y = 18
    max_chars = max(60, int((width - 42) / (font_size * 0.61)))
    for line in text.splitlines():
        if y + line_h > height - 12:
            break
        draw.text((22, y), line[:max_chars], font=mono, fill=(235, 235, 235))
        y += line_h
    img.save(out)
    return out


def file_text(path: Path, start: int = 1, end: int | None = None) -> str:
    lines = path.read_text(encoding="utf-8").splitlines()
    selected = lines[start - 1 : end]
    return "\n".join(f"{idx + start - 1:>3}  {line}" for idx, line in enumerate(selected))


def tree_text() -> str:
    proc = subprocess.run(
        ["find", str(BASE), "-maxdepth", "3", "-type", "f", "-printf", "%P\n"],
        text=True,
        stdout=subprocess.PIPE,
        check=True,
    )
    keep = []
    for line in sorted(proc.stdout.splitlines()):
        if (
            line.startswith((".git/", ".venv/"))
            or "__pycache__" in line
            or line.startswith("workspace/runs/")
            or line.startswith("outputs/")
            or line.startswith("presentation_assets/")
            or line.startswith("ppt_preview/")
            or line.endswith(".pptx")
            or line.endswith(".pdf")
        ):
            continue
        keep.append(line)
    return "\n".join(keep)


def add_footer(slide, slide_no: int, total: int) -> None:
    y = Inches(SLIDE_H - FOOTER_H)
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, y, Inches(SLIDE_W), Inches(FOOTER_H))
    footer.fill.solid()
    footer.fill.fore_color.rgb = BLUE
    footer.line.fill.background()

    gold = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(SLIDE_H - 0.035), Inches(SLIDE_W), Inches(0.035))
    gold.fill.solid()
    gold.fill.fore_color.rgb = GOLD
    gold.line.fill.background()

    # Simple UCI-like mark, drawn as clean white/teal wedges without using external branding files.
    x0 = Inches(0.24)
    y0 = Inches(SLIDE_H - FOOTER_H + 0.18)
    for i, col in enumerate([RGBColor(255, 255, 255), RGBColor(96, 188, 190), RGBColor(255, 255, 255)]):
        tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, x0 + Inches(i * 0.18), y0 + Inches(i * 0.03), Inches(0.36), Inches(0.28))
        tri.fill.solid()
        tri.fill.fore_color.rgb = col
        tri.line.fill.background()

    box = slide.shapes.add_textbox(Inches(2.55), Inches(SLIDE_H - FOOTER_H + 0.06), Inches(3.8), Inches(0.54))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "UCI Samueli"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = "School of Engineering"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(255, 255, 255)

    center = slide.shapes.add_textbox(Inches(7.0), Inches(SLIDE_H - FOOTER_H + 0.23), Inches(2.4), Inches(0.34))
    p = center.text_frame.paragraphs[0]
    p.text = "M.Eng Capstone Project"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    right = slide.shapes.add_textbox(Inches(11.65), Inches(SLIDE_H - FOOTER_H + 0.12), Inches(1.35), Inches(0.52))
    tf = right.text_frame
    p = tf.paragraphs[0]
    p.text = "Yunbo"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = f"{slide_no} / {total}"
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(255, 255, 255)


def add_title(slide, title: str, size: int = 30) -> None:
    box = slide.shapes.add_textbox(Inches(0.78), Inches(0.33), Inches(11.8), Inches(0.58))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(0, 0, 0)


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 20, color=RGBColor(0, 0, 0)) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_label_box(slide, text: str, x: float, y: float, w: float, h: float, fill, line, size: int = 16, bold: bool = False) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.4)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = DARK
    return shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color=RGBColor(30, 30, 30)) -> None:
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.8)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass


def add_ms_logo(slide) -> None:
    x = Inches(10.45)
    y = Inches(0.16)
    s = Inches(0.22)
    gap = Inches(0.035)
    colors = [MS_RED, MS_GREEN, MS_BLUE, MS_YELLOW]
    coords = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for (cx, cy), color in zip(coords, colors):
        sq = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + cx * (s + gap), y + cy * (s + gap), s, s)
        sq.fill.solid()
        sq.fill.fore_color.rgb = color
        sq.line.fill.background()
    box = slide.shapes.add_textbox(Inches(11.03), Inches(0.12), Inches(1.8), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = "Microsoft"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(110, 110, 110)


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_assets() -> dict[str, Path]:
    assets = {}
    assets["tree"] = code_image(tree_text(), ASSETS / "file_tree.png", width=980, font_size=20)
    assets["graph_code"] = code_panel(file_text(BASE / "codex_agent" / "graph.py", 20, 82), ASSETS / "graph_code.png")
    assets["node_code"] = code_panel(file_text(BASE / "codex_agent" / "nodes.py", 58, 126), ASSETS / "nodes_code.png")
    assets["tool_code"] = code_panel(file_text(BASE / "codex_agent" / "tools.py", 390, 468), ASSETS / "tools_code.png")
    cli = textwrap.dedent(
        """\
        cd MS_Migration_Agent
        python -m venv .venv
        source .venv/bin/activate
        pip install -r requirements.txt

        python run_agent.py \\
          --source examples/source_amplifier_example.scs \\
          --target-specs examples/target_specs_example.csv \\
          --target-pdk ptm22_lp \\
          --max-iterations 4 \\
          --prompt "Retarget the 45nm amplifier to PTM 22nm LP"

        Outputs:
          workspace/runs/<run_id>/
          outputs/<run_id>/final_target.scs
          outputs/<run_id>/measurements.ocn
          outputs/<run_id>/migration_report.md
        """
    )
    assets["cli"] = code_image(cli, ASSETS / "cli_demo.png", width=1120, font_size=21)
    return assets


def create_deck() -> Path:
    assets = build_assets()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    props = prs.core_properties
    props.title = "AURA — Agentic Retargeting Framework"
    props.subject = "Microsoft × UC Irvine M.Eng Capstone"
    props.author = "Yunbo Wang"
    props.last_modified_by = "Yunbo Wang"
    props.created = datetime(2026, 5, 8)
    props.modified = datetime(2026, 5, 8)
    total = 16

    # 1 Title
    slide = blank(prs)
    add_ms_logo(slide)
    for text, y, size, bold in [
        ("AURA", 0.68, 43, True),
        ("Agentic Retargeting Framework", 1.48, 34, True),
        ("LangGraph Implementation Update", 2.2, 24, True),
    ]:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(11.3), Inches(0.58))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(0, 0, 0)
    add_bullets(
        slide,
        [
            "Students: Yunbo Wang",
            "Company Liaisons: Joe Tostenrude, Yousef Iskander, Ph.D., Richard Paw",
            "Faculty Advisor: Dr. Farzad Ahmadkhanlou",
            "Date: May 08, 2026",
        ],
        2.05,
        3.52,
        9.2,
        1.5,
        size=16,
        color=RGBColor(25, 70, 120),
    )
    add_footer(slide, 1, total)

    # 2 Scope
    slide = blank(prs)
    add_title(slide, "This Deck: Agentic Framework Only")
    add_bullets(
        slide,
        [
            "Built the core LangGraph orchestration for analog netlist migration.",
            "Focus: agent roles, routing logic, Cadence/Ocean tool wrappers, run artifacts.",
            "Not covering model training, GNN decoder, or GLA experiments in this update.",
            "Default mode runs locally without UCI server; Cadence mode is enabled by environment variables.",
        ],
        1.0,
        1.4,
        5.8,
        4.1,
        size=22,
    )
    add_label_box(slide, "Input\n.scs + specs CSV + prompt", 7.4, 1.35, 2.0, 0.9, LIGHT_BLUE, UCI_BLUE, 15, True)
    add_label_box(slide, "LangGraph\nmulti-agent loop", 9.75, 1.35, 2.0, 0.9, LIGHT_GREEN, GREEN, 15, True)
    add_arrow(slide, 9.4, 1.8, 9.75, 1.8)
    add_label_box(slide, "Output\nfinal .scs + .ocn + report", 8.55, 3.05, 2.15, 0.95, LIGHT_ORANGE, ORANGE, 15, True)
    add_arrow(slide, 10.75, 2.25, 9.65, 3.05)
    add_footer(slide, 2, total)

    # 3 Existing flow
    slide = blank(prs)
    flow = BASE / "agent_flow.png"
    if flow.exists():
        slide.shapes.add_picture(str(flow), Inches(0.78), Inches(0.12), width=Inches(11.78))
    add_footer(slide, 3, total)

    # 4 LangGraph implementation
    slide = blank(prs)
    add_title(slide, "LangGraph Runtime")
    add_label_box(slide, "load_inputs", 0.75, 1.3, 1.55, 0.55, LIGHT_BLUE, UCI_BLUE, 12, True)
    add_label_box(slide, "retrieve_kb", 2.65, 1.3, 1.55, 0.55, LIGHT_BLUE, UCI_BLUE, 12, True)
    add_label_box(slide, "retarget_planner", 4.55, 1.3, 1.8, 0.55, LIGHT_GREEN, GREEN, 11, True)
    add_label_box(slide, "rule_engine", 6.75, 1.3, 1.55, 0.55, LIGHT_ORANGE, ORANGE, 12, True)
    add_label_box(slide, "spectre_compile", 8.65, 1.3, 1.8, 0.55, LIGHT_ORANGE, ORANGE, 11, True)
    add_label_box(slide, "measure / ocean", 10.9, 1.3, 1.75, 0.55, LIGHT_BLUE, UCI_BLUE, 12, True)
    for x in [2.3, 4.2, 6.35, 8.3, 10.45]:
        add_arrow(slide, x, 1.58, x + 0.42, 1.58)
    add_label_box(slide, "compile_debugger\nif Spectre errors", 8.42, 2.55, 1.8, 0.72, LIGHT_MAGENTA, MAGENTA, 12, True)
    add_arrow(slide, 9.2, 1.86, 9.2, 2.55, MAGENTA)
    add_arrow(slide, 8.42, 2.95, 7.25, 1.86, MAGENTA)
    add_label_box(slide, "performance_analyst", 5.15, 3.85, 1.9, 0.6, LIGHT_GREEN, GREEN, 13, True)
    add_label_box(slide, "optimizer", 7.8, 3.85, 1.3, 0.6, LIGHT_GREEN, GREEN, 13, True)
    add_label_box(slide, "finalize_outputs", 10.0, 3.85, 1.7, 0.6, LIGHT_BLUE, UCI_BLUE, 13, True)
    add_arrow(slide, 6.1, 1.85, 5.95, 3.85)
    add_arrow(slide, 7.05, 4.15, 7.8, 4.15)
    add_arrow(slide, 9.1, 4.15, 10.0, 4.15)
    add_arrow(slide, 8.45, 3.85, 6.9, 1.88, MAGENTA)
    add_footer(slide, 4, total)

    # 5 Graph code
    slide = blank(prs)
    add_title(slide, "Graph Code: Explicit Routing")
    slide.shapes.add_picture(str(assets["graph_code"]), Inches(0.65), Inches(1.05), width=Inches(12.0))
    add_footer(slide, 5, total)

    # 6 State contract
    slide = blank(prs)
    add_title(slide, "State Contract")
    columns = [
        ("Inputs", ["source_scs_path", "target_specs_csv_path", "target_pdk", "user_prompt"]),
        ("Retargeting", ["retarget_plan", "draft_scs_path", "rule_violations", "rule_pass"]),
        ("Tools", ["spectre_compile_ok", "ocean_script_path", "measured_csv_path", "compile_errors"]),
        ("Optimization", ["specs_met", "failed_specs", "optimization_actions", "final_scs_path"]),
    ]
    for i, (head, rows) in enumerate(columns):
        x = 0.75 + i * 3.1
        add_label_box(slide, head, x, 1.25, 2.35, 0.45, LIGHT_BLUE, UCI_BLUE, 16, True)
        add_bullets(slide, rows, x + 0.12, 1.92, 2.2, 2.7, size=17, color=DARK)
    add_bullets(
        slide,
        ["The graph passes one TypedDict state across every agent and tool node.", "All generated files are also registered in `artifacts` for reporting."],
        1.05,
        5.15,
        11.2,
        0.8,
        size=18,
        color=GRAY,
    )
    add_footer(slide, 6, total)

    # 7 Phase 1
    slide = blank(prs)
    add_title(slide, "Phase 1: Netlist Retargeting")
    add_bullets(
        slide,
        [
            "Reads source Spectre netlist and target specs CSV.",
            "Builds a first-pass retarget plan from local PDK presets and examples.",
            "Rewrites model include, MOS model names, default bias variables, and target L/W constraints.",
            "Rule engine checks model mapping, MOS pin count, min length/width, and required design parameters.",
        ],
        0.95,
        1.25,
        5.7,
        4.6,
        size=20,
    )
    add_label_box(slide, "45nm source\nnmos / pmos\n5V bias", 7.0, 1.35, 1.8, 0.95, LIGHT_BLUE, UCI_BLUE, 14, True)
    add_label_box(slide, "Target PDK preset\nptm22_lp / gpdk045 / asap7", 9.25, 1.35, 2.25, 0.95, LIGHT_ORANGE, ORANGE, 13, True)
    add_arrow(slide, 8.8, 1.82, 9.25, 1.82)
    add_label_box(slide, "Draft .scs\n22n L\nnew include\nsafe bias defaults", 8.05, 3.2, 2.35, 1.15, LIGHT_GREEN, GREEN, 13, True)
    add_arrow(slide, 10.35, 2.3, 9.22, 3.2)
    add_footer(slide, 7, total)

    # 8 Phase 1 code
    slide = blank(prs)
    add_title(slide, "Phase 1 Code: Planner and Rule Gate")
    slide.shapes.add_picture(str(assets["node_code"]), Inches(0.65), Inches(1.05), width=Inches(12.0))
    add_footer(slide, 8, total)

    # 9 Spectre compile
    slide = blank(prs)
    add_title(slide, "Tool Node: Spectre Compile / Debugger")
    add_bullets(
        slide,
        [
            "`run_spectre_compile` prepares the real Cadence command and parses logs.",
            "`AURA_DRY_RUN=true` skips external binaries for local development.",
            "When connected to UCI server, set `AURA_DRY_RUN=false` and source Cadence profile.",
            "Compile debugger keeps edits conservative: fix missing params, then route back to rule engine.",
        ],
        0.95,
        1.25,
        5.9,
        4.7,
        size=20,
    )
    cmd = "spectre -64 draft_target.scs +escchars +log spectre_compile.log +aps -maxw 5 -maxn 5"
    slide.shapes.add_picture(str(code_image(cmd, ASSETS / "spectre_cmd.png", width=950, font_size=22)), Inches(7.15), Inches(1.4), width=Inches(4.8))
    add_label_box(slide, "No SSH inside the agent\nRun locally on UCI Linux/Cadence host", 7.35, 3.0, 4.25, 0.9, LIGHT_MAGENTA, MAGENTA, 15, True)
    add_footer(slide, 9, total)

    # 10 Ocean
    slide = blank(prs)
    add_title(slide, "Phase 2: Topology-Aware Ocean Generation")
    add_bullets(
        slide,
        [
            "Infers differential amplifier probes: Vin+, Vin-, Vout1, Vout2, VDD, ground.",
            "Chooses analyses from requested metrics: dcOp, ac, transient.",
            "Generates `.ocn` with simulator, modelFile, desVar, analysis, resultsDir, ocnPrint.",
            "Normalizes Spectre values like `10uA`, `0.95V`, `10mV` into Ocean-friendly literals.",
        ],
        0.95,
        1.22,
        6.0,
        4.8,
        size=20,
    )
    add_label_box(slide, "Topology\nDetector", 7.35, 1.2, 1.55, 0.6, LIGHT_GREEN, GREEN, 14, True)
    add_label_box(slide, "Analysis\nPlanner", 9.2, 1.2, 1.55, 0.6, LIGHT_GREEN, GREEN, 14, True)
    add_label_box(slide, "Ocean\nWriter", 11.05, 1.2, 1.35, 0.6, LIGHT_BLUE, UCI_BLUE, 14, True)
    add_arrow(slide, 8.9, 1.5, 9.2, 1.5)
    add_arrow(slide, 10.75, 1.5, 11.05, 1.5)
    ocean_small = textwrap.dedent(
        """\
        simulator( 'spectre )
        modelFile( list( "22nm_LP.pm" "" ) )
        desVar( "VDD" 0.95 )
        analysis('tran ?stop "5u" )
        ocnPrint( v("/Vout2") ?output "waveforms.csv" )
        """
    )
    slide.shapes.add_picture(str(code_image(ocean_small, ASSETS / "ocean_small.png", width=800, font_size=22)), Inches(7.35), Inches(2.45), width=Inches(4.8))
    add_footer(slide, 10, total)

    # 11 Optimization
    slide = blank(prs)
    add_title(slide, "Phase 3: Performance Optimization Loop")
    add_label_box(slide, "measured_metrics.csv", 0.9, 1.35, 2.0, 0.55, LIGHT_BLUE, UCI_BLUE, 14, True)
    add_label_box(slide, "Performance Analyst\ncompare specs", 3.45, 1.2, 2.15, 0.85, LIGHT_GREEN, GREEN, 14, True)
    add_label_box(slide, "Specs Met?", 6.25, 1.25, 1.45, 0.72, LIGHT_MAGENTA, MAGENTA, 14, True)
    add_label_box(slide, "Final Report", 9.95, 1.35, 1.65, 0.55, LIGHT_BLUE, UCI_BLUE, 14, True)
    add_arrow(slide, 2.9, 1.62, 3.45, 1.62)
    add_arrow(slide, 5.6, 1.62, 6.25, 1.62)
    add_arrow(slide, 7.7, 1.62, 9.95, 1.62)
    add_label_box(slide, "Optimizer\nprioritize failed metrics", 5.05, 3.25, 2.45, 0.85, LIGHT_ORANGE, ORANGE, 14, True)
    add_label_box(slide, "Update .scs\nW/L, bias, R/C", 8.05, 3.25, 2.0, 0.85, LIGHT_GREEN, GREEN, 14, True)
    add_arrow(slide, 6.9, 1.98, 6.25, 3.25, MAGENTA)
    add_arrow(slide, 7.5, 3.68, 8.05, 3.68)
    add_arrow(slide, 9.05, 3.25, 6.95, 1.98, MAGENTA)
    add_bullets(
        slide,
        [
            "Gain or bandwidth failures tune Rg/Rm.",
            "Power failures reduce ISS/Ibtail.",
            "Output common-mode failures tune Aref/Vb3.",
        ],
        1.0,
        4.55,
        10.8,
        1.0,
        size=19,
        color=GRAY,
    )
    add_footer(slide, 11, total)

    # 12 Tool code
    slide = blank(prs)
    add_title(slide, "Tooling Code: Metrics and Optimization")
    slide.shapes.add_picture(str(assets["tool_code"]), Inches(0.65), Inches(1.05), width=Inches(12.0))
    add_footer(slide, 12, total)

    # 13 Code files
    slide = blank(prs)
    add_title(slide, "Agent Code Files")
    slide.shapes.add_picture(str(assets["tree"]), Inches(0.85), Inches(1.05), width=Inches(6.15))
    add_bullets(
        slide,
        [
            "- graph.py: LangGraph topology and conditional routers.",
            "- nodes.py: agent/tool nodes matching the three-stage flow.",
            "- tools.py: netlist parser, PDK presets, rule checks, runner wrappers.",
            "- state.py: shared state contract.",
            "- run_agent.py: reproducible CLI entrypoint.",
        ],
        7.35,
        1.35,
        4.8,
        3.8,
        size=18,
    )
    add_footer(slide, 13, total)

    # 14 CLI
    slide = blank(prs)
    add_title(slide, "How To Run The Agent")
    slide.shapes.add_picture(str(assets["cli"]), Inches(0.8), Inches(1.05), width=Inches(11.9))
    add_footer(slide, 14, total)

    # 15 Features
    slide = blank(prs)
    add_title(slide, "Current Technical Features")
    feature_boxes = [
        ("Local KB Retrieval", "Reads david / csc / TPM / mother_code examples for relevant snippets.", LIGHT_BLUE, UCI_BLUE),
        ("PDK Presets", "ptm22_lp, gpdk045, asap7 model names, min geometry, include paths.", LIGHT_GREEN, GREEN),
        ("Dry Run First", "Develop and demo without Cadence connection; no SSH required.", LIGHT_ORANGE, ORANGE),
        ("Cadence Ready", "Spectre/Ocean wrappers use environment variables for UCI server mode.", LIGHT_MAGENTA, MAGENTA),
        ("Artifact Tracking", "Every run writes draft .scs, .ocn, logs, metrics CSV, final report.", LIGHT_BLUE, UCI_BLUE),
        ("Loop Routing", "Optimizer sends updates back to rule checks or Ocean generation.", LIGHT_GREEN, GREEN),
    ]
    for i, (head, body, fill, line) in enumerate(feature_boxes):
        x = 0.8 + (i % 2) * 6.0
        y = 1.15 + (i // 2) * 1.45
        add_label_box(slide, head, x, y, 2.0, 0.5, fill, line, 14, True)
        add_bullets(slide, [body], x + 2.18, y + 0.05, 3.55, 0.62, size=15, color=DARK)
    add_footer(slide, 15, total)

    # 16 Next steps
    slide = blank(prs)
    add_title(slide, "Next Steps")
    add_bullets(
        slide,
        [
            "Run `AURA_DRY_RUN=false` on the UCI Cadence Linux host and validate Spectre/Ocean logs.",
            "Replace heuristic retarget planner with LLM + structured RAG over PDK manuals and tuning cases.",
            "Add stronger analog metric extraction: gain, bandwidth, phase margin, power, settling, swing.",
            "Connect Skillbridge/Virtuoso path for schematic-level desVar tuning.",
            "Later stage: log trajectories for Agent Lightning style policy/router improvement.",
        ],
        1.05,
        1.35,
        10.7,
        4.4,
        size=22,
    )
    add_footer(slide, 16, total)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = create_deck()
    print(path)
